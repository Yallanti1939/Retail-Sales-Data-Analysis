"""Run the full retail sales analysis and reporting pipeline."""

from __future__ import annotations

import argparse
from pathlib import Path
from shutil import copyfile

import pandas as pd
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils.dataframe import dataframe_to_rows
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table

from src.analysis import build_analysis_tables, calculate_kpis, generate_business_insights
from src.config import (
    BUSINESS_INSIGHTS_PATH,
    CHARTS_DIR,
    CLEANED_RETAIL_DATA_PATH,
    CLEANED_TRANSACTION_DATA_PATH,
    DATA_DIR,
    EXCEL_DASHBOARD_PATH,
    FINAL_REPORT_PATH,
    INSIGHTS_REPORT_PATH,
    PRESENTATION_PATH,
    PROJECT_ROOT,
    RAW_DATA_PATH,
    RAW_RETAIL_DATA_PATH,
)
from src.data_cleaning import clean_transactions
from src.data_loader import load_csv, profile_dataframe
from src.database import create_mysql_engine, initialize_database, insert_transactions
from src.feature_engineering import add_sales_features
from src.insights import export_business_insights
from src.visualization import create_all_charts


def save_cleaned_data(df: pd.DataFrame) -> Path:
    """Save the cleaned, feature-rich dataset under both requested names."""
    df.to_csv(CLEANED_TRANSACTION_DATA_PATH, index=False)
    df.to_csv(CLEANED_RETAIL_DATA_PATH, index=False)
    return CLEANED_RETAIL_DATA_PATH


def create_excel_dashboard(
    df: pd.DataFrame,
    tables: dict[str, pd.DataFrame],
    output_path: Path,
) -> Path:
    """Create an Excel dashboard workbook with KPIs, tables, and charts."""
    kpis = calculate_kpis(df)
    workbook = Workbook()
    dashboard = workbook.active
    dashboard.title = "Dashboard"

    header_fill = PatternFill("solid", fgColor="1F4E78")
    kpi_fill = PatternFill("solid", fgColor="D9EAF7")
    dashboard["A1"] = "Retail Sales Dashboard"
    dashboard["A1"].font = Font(size=18, bold=True, color="1F4E78")

    kpi_rows = [
        ("Total Revenue", kpis["total_revenue"]),
        ("Total Transactions", kpis["total_transactions"]),
        ("Total Countries", kpis["total_countries"]),
        ("Total Products", kpis["total_products"]),
        ("Average Order Value", kpis["average_order_value"]),
    ]
    for row_index, (label, value) in enumerate(kpi_rows, start=3):
        dashboard.cell(row=row_index, column=1, value=label)
        dashboard.cell(row=row_index, column=2, value=value)
        dashboard.cell(row=row_index, column=1).fill = kpi_fill
        dashboard.cell(row=row_index, column=1).font = Font(bold=True)
        dashboard.cell(row=row_index, column=2).number_format = "#,##0.00"

    for name, table_df in tables.items():
        worksheet = workbook.create_sheet(name[:31])
        for row in dataframe_to_rows(table_df, index=False, header=True):
            worksheet.append(row)
        for cell in worksheet[1]:
            cell.fill = header_fill
            cell.font = Font(bold=True, color="FFFFFF")
            cell.alignment = Alignment(horizontal="center")
        for column_cells in worksheet.columns:
            max_length = max(len(str(cell.value or "")) for cell in column_cells)
            worksheet.column_dimensions[column_cells[0].column_letter].width = min(
                max_length + 2,
                40,
            )

    monthly_sheet = workbook["monthly_sales"]
    line_chart = LineChart()
    line_chart.title = "Monthly Sales"
    line_chart.y_axis.title = "Revenue"
    line_chart.x_axis.title = "Month"
    data = Reference(monthly_sheet, min_col=2, min_row=1, max_row=monthly_sheet.max_row)
    cats = Reference(monthly_sheet, min_col=1, min_row=2, max_row=monthly_sheet.max_row)
    line_chart.add_data(data, titles_from_data=True)
    line_chart.set_categories(cats)
    dashboard.add_chart(line_chart, "D3")

    country_sheet = workbook["top_10_countries"]
    country_chart = BarChart()
    country_chart.title = "Top Countries"
    country_chart.y_axis.title = "Revenue"
    data = Reference(country_sheet, min_col=2, min_row=1, max_row=country_sheet.max_row)
    cats = Reference(country_sheet, min_col=1, min_row=2, max_row=country_sheet.max_row)
    country_chart.add_data(data, titles_from_data=True)
    country_chart.set_categories(cats)
    dashboard.add_chart(country_chart, "D19")

    product_sheet = workbook["top_10_products"]
    product_chart = BarChart()
    product_chart.title = "Top Products"
    data = Reference(product_sheet, min_col=3, min_row=1, max_row=product_sheet.max_row)
    cats = Reference(product_sheet, min_col=2, min_row=2, max_row=product_sheet.max_row)
    product_chart.add_data(data, titles_from_data=True)
    product_chart.set_categories(cats)
    dashboard.add_chart(product_chart, "L3")

    hourly_sheet = workbook["hourly_transactions"]
    hourly_chart = BarChart()
    hourly_chart.title = "Peak Hours"
    data = Reference(hourly_sheet, min_col=2, min_row=1, max_row=hourly_sheet.max_row)
    cats = Reference(hourly_sheet, min_col=1, min_row=2, max_row=hourly_sheet.max_row)
    hourly_chart.add_data(data, titles_from_data=True)
    hourly_chart.set_categories(cats)
    dashboard.add_chart(hourly_chart, "L19")

    dashboard.column_dimensions["A"].width = 24
    dashboard.column_dimensions["B"].width = 18
    output_path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output_path)
    return output_path


def create_pdf_report(
    kpis: dict[str, float | int],
    insights: list[str],
    chart_paths: list[Path],
    output_path: Path,
) -> Path:
    """Create a concise PDF project report."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    document = SimpleDocTemplate(str(output_path), pagesize=letter)
    elements = [
        Paragraph("Retail Sales Analysis and Reporting", styles["Title"]),
        Spacer(1, 12),
        Paragraph("Executive KPI Summary", styles["Heading2"]),
    ]
    kpi_table = Table(
        [["Metric", "Value"]]
        + [[key.replace("_", " ").title(), f"{value:,.2f}"] for key, value in kpis.items()]
    )
    kpi_table.setStyle(
        [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E78")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ]
    )
    elements.extend([kpi_table, Spacer(1, 12), Paragraph("Business Insights", styles["Heading2"])])
    for insight in insights:
        elements.append(Paragraph(f"- {insight}", styles["BodyText"]))
    elements.append(Spacer(1, 12))
    elements.append(Paragraph("Selected Visualizations", styles["Heading2"]))
    for chart_path in chart_paths[:6]:
        elements.append(Image(str(chart_path), width=480, height=260))
        elements.append(Spacer(1, 10))
    document.build(elements)
    return output_path


def create_presentation(
    kpis: dict[str, float | int],
    insights: list[str],
    chart_paths: list[Path],
    output_path: Path,
) -> Path:
    """Create a PowerPoint presentation summarizing the project."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Retail Sales Analysis"
    title_slide.placeholders[1].text = "End-to-end Python, SQL, Excel, and dashboard project"

    kpi_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    kpi_slide.shapes.title.text = "KPI Summary"
    left = Inches(0.7)
    top = Inches(1.4)
    for index, (label, value) in enumerate(kpis.items()):
        textbox = kpi_slide.shapes.add_textbox(left, top + Inches(index * 0.7), Inches(8.6), Inches(0.45))
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.text = f"{label.replace('_', ' ').title()}: {value:,.2f}"
        paragraph.font.size = Pt(20)

    charts_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    charts_slide.shapes.title.text = "Sales Trends and Revenue Drivers"
    for index, chart_path in enumerate(chart_paths[:2]):
        charts_slide.shapes.add_picture(
            str(chart_path),
            Inches(0.6 + index * 4.7),
            Inches(1.3),
            width=Inches(4.3),
        )

    insight_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    insight_slide.shapes.title.text = "Business Insights"
    body = insight_slide.placeholders[1].text_frame
    body.clear()
    for insight in insights[:6]:
        paragraph = body.add_paragraph()
        paragraph.text = insight
        paragraph.level = 0
        paragraph.font.size = Pt(18)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def run_pipeline(load_mysql: bool = False) -> dict[str, Path]:
    """Execute the full analysis pipeline and return generated artifact paths."""
    raw_path = RAW_RETAIL_DATA_PATH if RAW_RETAIL_DATA_PATH.exists() else RAW_DATA_PATH
    raw_df = load_csv(raw_path)
    profile = profile_dataframe(raw_df)
    print(f"Raw shape: {profile['shape']}")
    print(raw_df.head())
    print(raw_df.tail())
    print(raw_df.describe(include='all'))

    cleaned_df = clean_transactions(raw_df)
    featured_df = add_sales_features(cleaned_df)
    cleaned_path = save_cleaned_data(featured_df)

    if load_mysql:
        engine = create_mysql_engine()
        initialize_database(engine)
        insert_transactions(featured_df, engine)

    tables = build_analysis_tables(featured_df)
    kpis = calculate_kpis(featured_df)
    insights = generate_business_insights(featured_df)
    insights_path = export_business_insights(featured_df, BUSINESS_INSIGHTS_PATH)
    chart_paths = create_all_charts(featured_df, CHARTS_DIR)

    excel_path = create_excel_dashboard(
        featured_df,
        tables,
        EXCEL_DASHBOARD_PATH,
    )
    insights_pdf_path = create_pdf_report(
        kpis,
        insights,
        chart_paths,
        INSIGHTS_REPORT_PATH,
    )
    pdf_path = FINAL_REPORT_PATH
    copyfile(insights_pdf_path, pdf_path)
    pptx_path = create_presentation(
        kpis,
        insights,
        chart_paths,
        PRESENTATION_PATH,
    )

    return {
        "cleaned_data": cleaned_path,
        "business_insights": insights_path,
        "insights_report": insights_pdf_path,
        "excel_dashboard": excel_path,
        "pdf_report": pdf_path,
        "presentation": pptx_path,
    }


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Run the retail sales analysis and reporting pipeline."
    )
    parser.add_argument(
        "--load-mysql",
        action="store_true",
        help=(
            "Create the MySQL transactions table and insert cleaned data. "
            "Uses MYSQL_HOST, MYSQL_PORT, MYSQL_USER, MYSQL_PASSWORD, "
            "and MYSQL_DATABASE environment variables."
        ),
    )
    args = parser.parse_args()
    outputs = run_pipeline(load_mysql=args.load_mysql)
    for artifact, path in outputs.items():
        print(f"{artifact}: {path}")
